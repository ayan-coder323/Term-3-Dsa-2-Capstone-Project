from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG_DARK      = RGBColor(0x0F, 0x1B, 0x2D)
BG_CARD      = RGBColor(0x16, 0x2A, 0x45)
ACCENT       = RGBColor(0x38, 0xB2, 0xAC)
ACCENT_LIGHT = RGBColor(0x81, 0xE6, 0xD9)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE0)
MUTED        = RGBColor(0xA0, 0xAE, 0xC0)
GOLD         = RGBColor(0xED, 0xD8, 0x8B)
ORANGE       = RGBColor(0xF6, 0xAD, 0x55)
GREEN        = RGBColor(0x68, 0xD3, 0x91)
PINK         = RGBColor(0xED, 0x64, 0xA6)
PURPLE       = RGBColor(0xB7, 0x94, 0xF6)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def connect_shapes(slide, start_shape, end_shape):
    # This is a bit tricky in pptx, we'll just draw a manual arrow
    pass

def draw_arrow(slide, start_x, start_y, end_x, end_y):
    # Draw simple connector
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = MUTED
    connector.line.width = Pt(2)
    return connector

def update_problem_statement(prs):
    slide = prs.slides[1] # SLIDE 2
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "Organizations and students struggle" in shape.text or "manage documents efficiently" in shape.text:
                new_text = ("Modern organizations and academic institutions face significant challenges in managing large volumes of documents efficiently. Traditional file systems lack the structured indexing required for rapid retrieval, the relational mapping necessary to identify document dependencies, and the analytical tools needed for deep insights.\n\n"
                            "DocuTrack addresses this gap by implementing a robust, unified Data Structures & Algorithms (DSA) framework. By leveraging self-balancing trees, graph networks, and optimization algorithms, it transforms a flat document repository into a highly queryable, scalable, and interconnected ecosystem, providing instantaneous search, smart relationships, and intelligent storage allocation.")
                shape.text = new_text
                tf = shape.text_frame
                for p in tf.paragraphs:
                    p.font.size = Pt(15)
                    p.font.color.rgb = LIGHT_GRAY
                    p.font.name = "Segoe UI"
                    p.space_after = Pt(10)
                    p.line_spacing = 1.2
                break

def insert_menu_flow_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "System Flow & Menu Architecture", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

    # Flow blocks
    # 1. Start Box
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.5), Inches(3.1), Inches(0.7), BG_CARD, ACCENT)
    add_text_box(slide, Inches(5.1), Inches(1.65), Inches(3.1), Inches(0.4), "Start DocuTrack", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Arrow down
    draw_arrow(slide, Inches(6.65), Inches(2.2), Inches(6.65), Inches(2.6))

    # 2. Authentication
    add_shape(slide, MSO_SHAPE.HEXAGON, Inches(5.1), Inches(2.6), Inches(3.1), Inches(0.8), BG_CARD, GOLD)
    add_text_box(slide, Inches(5.1), Inches(2.8), Inches(3.1), Inches(0.4), "Authentication", font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Branches
    draw_arrow(slide, Inches(5.1), Inches(3.0), Inches(3.0), Inches(3.0)) # left
    draw_arrow(slide, Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.6))

    draw_arrow(slide, Inches(8.2), Inches(3.0), Inches(10.3), Inches(3.0)) # right
    draw_arrow(slide, Inches(10.3), Inches(3.0), Inches(10.3), Inches(3.6))

    # 3. Admin Branch
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.6), Inches(3.0), Inches(0.6), BG_CARD, ACCENT_LIGHT)
    add_text_box(slide, Inches(1.5), Inches(3.7), Inches(3.0), Inches(0.4), "Admin Login", font_size=16, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

    draw_arrow(slide, Inches(3.0), Inches(4.2), Inches(3.0), Inches(4.6))

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.6), Inches(4.0), Inches(2.5), BG_DARK, ACCENT_LIGHT)
    add_text_box(slide, Inches(1.0), Inches(4.7), Inches(4.0), Inches(0.4), "Admin Dashboard (Full Access)", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    modules = ["1. Document Indexing (BST/AVL)", "2. Document Analytics (Fenwick/B-Tree)", "3. Document Network (Graph)", "4. Path Optimization (Dijkstra)", "5. Sorting & Ranking", "6. Storage Optimization", "7. Exit"]
    for i, mod in enumerate(modules):
        add_text_box(slide, Inches(1.2), Inches(5.1 + i*0.25), Inches(3.6), Inches(0.3), mod, font_size=12, color=LIGHT_GRAY)

    # 4. User Branch
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(3.6), Inches(3.0), Inches(0.6), BG_CARD, GREEN)
    add_text_box(slide, Inches(8.8), Inches(3.7), Inches(3.0), Inches(0.4), "User Login", font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    draw_arrow(slide, Inches(10.3), Inches(4.2), Inches(10.3), Inches(4.6))

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(4.6), Inches(4.0), Inches(1.5), BG_DARK, GREEN)
    add_text_box(slide, Inches(8.3), Inches(4.7), Inches(4.0), Inches(0.4), "User Dashboard (Read-Only)", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    user_modules = ["1. View All Documents", "2. Search Document", "3. Exit"]
    for i, mod in enumerate(user_modules):
        add_text_box(slide, Inches(8.5), Inches(5.1 + i*0.25), Inches(3.6), Inches(0.3), mod, font_size=12, color=LIGHT_GRAY)

    # Move slide to index 2
    slide_id = slide.slide_id
    slide_list = prs.slides._sldIdLst
    for el in slide_list:
        if el.id == slide_id:
            new_el = el
            break
    slide_list.insert(2, new_el)

def insert_class_diagram_slide(prs):
    # Find and delete the existing class diagram slide (it's Slide 6, index 5 if we haven't inserted Menu flow yet, but we will have inserted Menu flow, so it becomes index 6)
    # Wait, the old Class Diagram slide is originally Slide 6 (index 5 in python-pptx before my insertions).
    # After inserting Menu Flow at index 2, the old Class Diagram becomes index 6.
    
    # First, let's just delete the old picture slide
    slide_list = prs.slides._sldIdLst
    old_class_diag_idx = 6 # because we inserted one before it
    slide_list.remove(slide_list[old_class_diag_idx])
    
    # Now create the new one
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "Class Diagram & Architecture", font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

    # Let's draw some UML style boxes
    def draw_class(slide, left, top, width, name, fields, methods, color=ACCENT):
        # Calculate height based on lines
        lines = 1 + len(fields) + len(methods)
        height = Inches(0.5 + lines * 0.2)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, BG_CARD, color)
        
        add_text_box(slide, left, top + Inches(0.05), width, Inches(0.3), name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        
        # separator
        draw_arrow(slide, left, top + Inches(0.4), left + width, top + Inches(0.4))
        
        y = top + Inches(0.45)
        for f in fields:
            add_text_box(slide, left + Inches(0.1), y, width - Inches(0.2), Inches(0.2), f, font_size=10, color=LIGHT_GRAY)
            y += Inches(0.2)
            
        draw_arrow(slide, left, y, left + width, y)
        y += Inches(0.05)
        
        for m in methods:
            add_text_box(slide, left + Inches(0.1), y, width - Inches(0.2), Inches(0.2), m, font_size=10, color=WHITE)
            y += Inches(0.2)

    # 1. Main Controller
    draw_class(slide, Inches(4.5), Inches(1.5), Inches(4.3), "Main (Controller)", 
               ["- documents: ArrayList<Document>", "- bst: BST", "- avl: AVLTree", "- graph: Graph"], 
               ["+ main(String[] args)", "+ login()", "+ displayAdminMenu()", "+ rebuildIndexes()"], PINK)

    # 2. Document Entity
    draw_class(slide, Inches(1.0), Inches(1.5), Inches(2.8), "Document", 
               ["- documentId: int", "- title: String", "- author: String"], 
               ["+ getters()", "+ setters()", "+ toString()"], GOLD)

    # 3. Indexing
    draw_class(slide, Inches(1.0), Inches(4.2), Inches(2.8), "BST / AVLTree", 
               ["- root: Node"], 
               ["+ insert(Document)", "+ search(int id)", "+ inorder()"], ACCENT)

    # 4. Graph
    draw_class(slide, Inches(4.5), Inches(4.2), Inches(4.3), "Graph", 
               ["- adjList: LinkedList<Edge>[]"], 
               ["+ addEdge(src, dest, weight)", "+ bfs()", "+ dfs()", "+ dijkstra()"], GREEN)

    # 5. Sorting
    draw_class(slide, Inches(9.5), Inches(1.5), Inches(3.0), "SortingManager", 
               ["- data: List<Document>"], 
               ["+ mergeSort()", "+ quickSort()", "+ heapSort()"], PURPLE)

    # 6. Optimization
    draw_class(slide, Inches(9.5), Inches(4.2), Inches(3.0), "OptimizationManager", 
               ["- items: List<Document>"], 
               ["+ knapsack01()", "+ activitySelection()", "+ longestIncSubseq()"], ORANGE)

    # Draw connectors (just dashed lines indicating relationship)
    # Document <- Main
    c1 = draw_arrow(slide, Inches(3.8), Inches(2.5), Inches(4.5), Inches(2.5))
    # Indexing <- Main
    c2 = draw_arrow(slide, Inches(2.4), Inches(4.2), Inches(4.5), Inches(3.2))
    # Graph <- Main
    c3 = draw_arrow(slide, Inches(6.6), Inches(3.5), Inches(6.6), Inches(4.2))
    # Sorting <- Main
    c4 = draw_arrow(slide, Inches(8.8), Inches(2.5), Inches(9.5), Inches(2.5))
    # Optimization <- Main
    c5 = draw_arrow(slide, Inches(8.8), Inches(3.2), Inches(9.5), Inches(4.5))

    # Move to index 6
    slide_id = slide.slide_id
    slide_list = prs.slides._sldIdLst
    for el in slide_list:
        if el.id == slide_id:
            new_el = el
            break
    slide_list.insert(6, new_el)

def main():
    prs = Presentation("DocuTrack_Presentation_final.pptx")
    update_problem_statement(prs)
    insert_menu_flow_slide(prs)
    insert_class_diagram_slide(prs)
    prs.save("DocuTrack_Presentation_final_Upgraded.pptx")
    print("Done")

if __name__ == "__main__":
    main()
