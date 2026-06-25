"""
DocuTrack – Professional Presentation Generator
Generates a soothing-themed PowerPoint with all project details.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette (Soft teal / navy / cream) ──────────────────────────
BG_DARK      = RGBColor(0x0F, 0x1B, 0x2D)   # deep navy
BG_CARD      = RGBColor(0x16, 0x2A, 0x45)   # card navy
ACCENT       = RGBColor(0x38, 0xB2, 0xAC)   # teal
ACCENT_LIGHT = RGBColor(0x81, 0xE6, 0xD9)   # light teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE0)
MUTED        = RGBColor(0xA0, 0xAE, 0xC0)
GOLD         = RGBColor(0xED, 0xD8, 0x8B)
ORANGE       = RGBColor(0xF6, 0xAD, 0x55)
PINK         = RGBColor(0xED, 0x64, 0xA6)
PURPLE       = RGBColor(0xB7, 0x94, 0xF6)
GREEN        = RGBColor(0x68, 0xD3, 0x91)
TABLE_HDR    = RGBColor(0x1A, 0x36, 0x5D)
TABLE_ROW1   = RGBColor(0x1E, 0x3A, 0x5F)
TABLE_ROW2   = RGBColor(0x16, 0x2A, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helper functions ────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, rows, cols, col_widths, data, font_size=11):
    """data = list of lists; first row is header."""
    table_width = sum(col_widths)
    row_height = Inches(0.42)
    total_height = row_height * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, total_height)
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = str(data[ri][ci])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Segoe UI"
                p.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_LIGHT
                else:
                    p.font.color.rgb = LIGHT_GRAY

            fill = cell.fill
            fill.solid()
            if ri == 0:
                fill.fore_color.rgb = TABLE_HDR
            elif ri % 2 == 1:
                fill.fore_color.rgb = TABLE_ROW1
            else:
                fill.fore_color.rgb = TABLE_ROW2
    return table_shape


def section_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)

    # Decorative left accent bar
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

    # Title
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5), title,
                 font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    add_accent_line(slide, Inches(1), Inches(3.7), Inches(3), ACCENT)

    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1), subtitle,
                     font_size=20, color=MUTED, alignment=PP_ALIGN.LEFT)
    return slide


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Top accent line
add_shape(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)

# Icon-like circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(1.2), Inches(1.8), Inches(1.8))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT
circle.line.fill.background()
add_text_box(slide, Inches(5.7), Inches(1.55), Inches(1.8), Inches(1.1), "📄", font_size=48, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1), "DocuTrack",
             font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI")

add_accent_line(slide, Inches(5.2), Inches(4.3), Inches(2.8), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.7),
             "Document Tracking & Management System",
             font_size=22, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.6),
             "A comprehensive DSA-powered console application in Java",
             font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.1), W, Inches(0.4), BG_CARD)
add_text_box(slide, Inches(1), Inches(7.12), Inches(11), Inches(0.35),
             "Java SE 21  •  Console Application  •  6 CO Modules  •  20+ DSA Implementations",
             font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: Project Overview
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "Project Overview",
             font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8),
             "DocuTrack is a Java SE 21 console application that demonstrates the practical application of "
             "Data Structures & Algorithms to a real-world document management problem.",
             font_size=15, color=LIGHT_GRAY)

# CO overview table
data = [
    ["CO", "Focus Area", "Data Structures & Algorithms"],
    ["CO1", "Document Indexing", "BST, AVL Tree"],
    ["CO2", "Analytics", "Fenwick Tree, Segment Tree, B-Tree, B+ Tree"],
    ["CO3", "Document Network", "Graph (Adjacency List), BFS, DFS, MST"],
    ["CO4", "Path Optimization", "Dijkstra, Bellman-Ford, Floyd-Warshall, Topological Sort"],
    ["CO5", "Sorting & Ranking", "Merge Sort, Quick Sort, Heap Sort, Radix Sort"],
    ["CO6", "Storage Optimization", "Activity Selection, Fractional Knapsack, 0/1 Knapsack, LIS"],
]
add_table(slide, Inches(0.8), Inches(2.5), 7, 3,
          [Inches(1.2), Inches(2.5), Inches(8)], data, font_size=12)

# Key highlight box
add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7), BG_CARD, ACCENT)
add_text_box(slide, Inches(1.1), Inches(6.05), Inches(11), Inches(0.6),
             "🔑  All modules share a single, unified document pool — inserting in BST makes data available to Graph, Path, Analytics, and all COs.",
             font_size=13, color=ACCENT_LIGHT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: Project Structure
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "Project Structure & Data Model",
             font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2.5))

# File tree
tree_text = (
    "DocuTrack/src/\n"
    "├── app/Main.java                  ← Entry point & menus\n"
    "├── models/Document.java           ← Core data model\n"
    "└── modules/\n"
    "    ├── indexing/bst/  (BST, BSTNode)\n"
    "    ├── indexing/avl/  (AVLTree, AVLNode)\n"
    "    ├── analytics/     (Fenwick, Segment, BTree, B+Tree)\n"
    "    ├── graph/         (Graph, BFS, DFS, MST)\n"
    "    ├── path/          (Dijkstra, BellmanFord, FloydWarshall, TopologicalSort)\n"
    "    ├── sorting/       (MergeSort, QuickSort, HeapSort, RadixSort)\n"
    "    └── optimization/  (ActivitySelection, FractionalKnapsack, Knapsack01, LIS)"
)
add_shape(slide, Inches(0.6), Inches(1.4), Inches(7.5), Inches(4.2), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(7), Inches(4), tree_text,
             font_size=13, color=GREEN, font_name="Consolas")

# Document model card
add_shape(slide, Inches(8.5), Inches(1.4), Inches(4.3), Inches(4.2), BG_CARD, ACCENT)
add_text_box(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(0.5), "models.Document",
             font_size=20, color=ACCENT_LIGHT, bold=True)

doc_fields = [
    "Fields:",
    "  • documentId : int",
    "  • title : String",
    "  • author : String",
    "",
    "Methods:",
    "  • Document(id, title, author)",
    "  • toString()",
    "",
    "Output Format:",
    '  "ID: 101, Title: Report Q1,',
    '   Author: Alice"',
]
add_bullet_list(slide, Inches(8.8), Inches(2.2), Inches(3.8), Inches(3.2), doc_fields, font_size=12, color=LIGHT_GRAY)

# Auth card
add_shape(slide, Inches(0.6), Inches(5.9), Inches(12.2), Inches(1.2), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(5.95), Inches(3), Inches(0.4), "Authentication & Roles",
             font_size=18, color=GOLD, bold=True)
add_text_box(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.7),
             "Admin (admin / admin123) → Full access to all 6 CO modules + document CRUD    |    "
             "User (user / user123) → Read-only: Search & View documents only",
             font_size=13, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: CO1 – Document Indexing
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.5), "CO1",
             font_size=16, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), "Document Indexing — BST & AVL Tree",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5))

# BST card
add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5), Inches(0.5), "Binary Search Tree (BST)",
             font_size=20, color=ACCENT_LIGHT, bold=True)
bst_items = [
    "Classes: BST.java, BSTNode.java",
    "Insert — O(h) iterative walk-down",
    "Search — O(h) iterative lookup",
    "Delete — 3 cases: leaf, one child, two children",
    "  → Two children uses in-order predecessor",
    "Traversals: Inorder, Preorder, Postorder",
    "Min / Max — leftmost / rightmost node",
    "",
    "Menu: Insert, Search, Delete, Inorder,",
    "         Minimum, Maximum, Back"
]
add_bullet_list(slide, Inches(0.9), Inches(2.3), Inches(5.3), Inches(4.3), bst_items, font_size=12)

# AVL card
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.3), BG_CARD, PURPLE)
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "AVL Tree (Self-Balancing BST)",
             font_size=20, color=PURPLE, bold=True)
avl_items = [
    "Classes: AVLTree.java, AVLNode.java",
    "Insert — O(log n) recursive + rotations",
    "Search — O(log n) recursive lookup",
    "Guarantees balanced tree (height tracked)",
    "",
    "Rotation Cases:",
    "  LL → Right Rotate",
    "  RR → Left Rotate",
    "  LR → Left Rotate child, then Right Rotate",
    "  RL → Right Rotate child, then Left Rotate",
    "",
    "Menu: Insert, Search, Inorder, Back"
]
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.3), Inches(4.3), avl_items, font_size=12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: CO2 – Analytics
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.5), "CO2",
             font_size=16, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), "Analytics — Fenwick, Segment, B-Tree, B+ Tree",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5))

# 4 cards in 2x2
cards = [
    ("Fenwick Tree (BIT)", ACCENT, [
        "Point update + Prefix sum",
        "update(i, v) — O(log n)",
        "sum(i) — O(log n)",
        "rangeSum(l, r)",
        "Use: Cumulative doc access counts"
    ]),
    ("Segment Tree", GOLD, [
        "Range sum queries — O(log n)",
        "build() — O(n)",
        "query(node, s, e, l, r)",
        "Use: Sum of doc IDs in a range",
        "e.g. \"Total IDs from index 2–5\""
    ]),
    ("B-Tree", GREEN, [
        "Multi-level key storage",
        "insert(key) — recursive",
        "display() — inorder traversal",
        "Stores document IDs as keys",
        "Use: Demonstrate multi-level indexing"
    ]),
    ("B+ Tree", ORANGE, [
        "Range search over doc IDs",
        "insert(key)",
        "rangeSearch(start, end)",
        "Use: Find all doc IDs in [100, 200]",
        "Sequential scan of stored data"
    ]),
]

positions = [
    (Inches(0.6), Inches(1.6)),
    (Inches(6.8), Inches(1.6)),
    (Inches(0.6), Inches(4.3)),
    (Inches(6.8), Inches(4.3)),
]

for (title, color, items), (l, t) in zip(cards, positions):
    add_shape(slide, l, t, Inches(5.8), Inches(2.4), BG_CARD, color)
    add_text_box(slide, l + Inches(0.3), t + Inches(0.1), Inches(5), Inches(0.4), title,
                 font_size=18, color=color, bold=True)
    add_bullet_list(slide, l + Inches(0.3), t + Inches(0.55), Inches(5.2), Inches(1.7), items, font_size=12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: CO3 – Document Network
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.5), "CO3",
             font_size=16, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), "Document Network — Graph, BFS, DFS, MST",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5))

# Graph card
add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5), Inches(0.5), "Graph (Adjacency List)",
             font_size=20, color=ACCENT_LIGHT, bold=True)
graph_items = [
    "Undirected graph — each vertex = document",
    "addEdge(src, dst) — bidirectional",
    "resize(n) — grows/shrinks with doc pool",
    "removeVertex(v) — removes + re-indexes",
    "display() — prints adjacency list",
    "displayDocuments() — with titles",
    "",
    "Menu: View Docs, Add Relationship,",
    "  Display Graph, BFS, DFS, MST, Back"
]
add_bullet_list(slide, Inches(0.9), Inches(2.3), Inches(5.3), Inches(4.3), graph_items, font_size=12)

# BFS/DFS/MST card
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.3), BG_CARD, GREEN)
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5), "Traversals & MST",
             font_size=20, color=GREEN, bold=True)
traversal_items = [
    "BFS (Breadth-First Search):",
    "  • Queue-based level-order traversal",
    "  • Visits all reachable vertices",
    "  • Complexity: O(V + E)",
    "",
    "DFS (Depth-First Search):",
    "  • Recursive depth exploration",
    "  • Uses boolean[] visited array",
    "  • Complexity: O(V + E)",
    "",
    "MST (Minimum Spanning Tree):",
    "  • Placeholder demonstration",
    "  • Infrastructure ready for Prim/Kruskal"
]
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.3), Inches(4.3), traversal_items, font_size=12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: CO4 – Path Optimization
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.5), "CO4",
             font_size=16, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), "Path Optimization — Shortest Path Algorithms",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5))

path_cards = [
    ("Dijkstra", ACCENT, [
        "Single-source shortest path (greedy)",
        "Input: cost[][], source, destination",
        "Output: min cost + full path",
        "Tracks parent[] for path reconstruction",
        "Complexity: O(V²)"
    ]),
    ("Bellman-Ford", GOLD, [
        "Single-source (handles negatives)",
        "Input: cost[][], source",
        "Output: distances to ALL vertices",
        "V-1 relaxation iterations",
        "Complexity: O(V³) with matrix"
    ]),
    ("Floyd-Warshall", PINK, [
        "All-pairs shortest path (DP)",
        "Input: cost[][]",
        "Output: full N×N distance matrix",
        "Triple nested loop via intermediate k",
        "Complexity: O(V³)"
    ]),
    ("Topological Sort", GREEN, [
        "DFS-based linear ordering (DAG)",
        "Input: cost[][] as directed graph",
        "Output: ordered vertex sequence",
        "Uses Stack for reverse post-order",
        "Complexity: O(V + E)"
    ]),
]

for (title, color, items), (l, t) in zip(path_cards, positions):
    add_shape(slide, l, t, Inches(5.8), Inches(2.4), BG_CARD, color)
    add_text_box(slide, l + Inches(0.3), t + Inches(0.1), Inches(5), Inches(0.4), title,
                 font_size=18, color=color, bold=True)
    add_bullet_list(slide, l + Inches(0.3), t + Inches(0.55), Inches(5.2), Inches(1.7), items, font_size=12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: CO5 – Sorting & Ranking
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.5), "CO5",
             font_size=16, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), "Sorting & Ranking — 4 Sorting Algorithms",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5))

# Table
sort_data = [
    ["Algorithm", "Sorts By", "Stable?", "Avg Time", "Technique"],
    ["Merge Sort", "Document ID", "✅ Yes", "O(n log n)", "Divide & Conquer"],
    ["Quick Sort", "Title (A-Z)", "❌ No", "O(n log n)", "Partition-based"],
    ["Heap Sort", "Author (A-Z)", "❌ No", "O(n log n)", "Max-Heap"],
    ["Radix Sort", "Document ID", "✅ Yes", "O(d × n)", "Non-comparative (LSD)"],
]
add_table(slide, Inches(0.8), Inches(1.6), 5, 5,
          [Inches(2.5), Inches(2.5), Inches(1.5), Inches(2.5), Inches(3)], sort_data, font_size=13)

# Details
add_shape(slide, Inches(0.6), Inches(4.0), Inches(12.2), Inches(3.1), BG_CARD)
sort_details = [
    "All algorithms operate on a COPY of the document list — original data is never modified.",
    "Merge Sort — Classic divide-and-conquer: split, recurse, merge by documentId comparison.",
    "Quick Sort — Lomuto partition using title.toLowerCase().compareTo() as comparator.",
    "Heap Sort — Builds max-heap using author.compareToIgnoreCase(), then extracts max iteratively.",
    "Radix Sort — LSD (Least Significant Digit) approach with counting sort per digit position.",
]
add_bullet_list(slide, Inches(0.9), Inches(4.1), Inches(11.5), Inches(2.8), sort_details, font_size=13)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: CO6 – Storage Optimization
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(2), Inches(0.5), "CO6",
             font_size=16, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), "Storage Optimization — Greedy & DP",
             font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5))

opt_cards = [
    ("Activity Selection (Greedy)", ACCENT, [
        "Input: documents + start/finish times",
        "Selects non-overlapping activities",
        "Sorted by finish time (greedy choice)",
        "Use: Schedule doc review sessions"
    ]),
    ("Fractional Knapsack (Greedy)", GOLD, [
        "Input: docs + value/weight + capacity",
        "Picks by value-to-weight ratio",
        "Can take fractional portions",
        "Use: Partial storage allocation"
    ]),
    ("0/1 Knapsack (DP)", PINK, [
        "Input: docs + value/weight + capacity",
        "DP table: dp[i][w] = max value",
        "Backtrack to find selected docs",
        "Use: Optimal subset selection"
    ]),
    ("LIS — Longest Increasing Subseq (DP)", GREEN, [
        "Input: array of document IDs",
        "Classic O(n²) DP approach",
        "Tracks parent[] for reconstruction",
        "Use: Find increasing ID pattern"
    ]),
]

for (title, color, items), (l, t) in zip(opt_cards, positions):
    add_shape(slide, l, t, Inches(5.8), Inches(2.4), BG_CARD, color)
    add_text_box(slide, l + Inches(0.3), t + Inches(0.1), Inches(5), Inches(0.4), title,
                 font_size=16, color=color, bold=True)
    add_bullet_list(slide, l + Inches(0.3), t + Inches(0.55), Inches(5.2), Inches(1.7), items, font_size=12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: Execution Flow
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Execution Flow",
             font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# Flow boxes
flow_steps = [
    ("LOGIN", "admin/admin123\nuser/user123", ACCENT, Inches(0.8)),
    ("ROLE\nCHECK", "Admin → Full\nUser → View Only", GOLD, Inches(3.3)),
    ("ADMIN\nMENU", "CO1–CO6\nSub-menus", GREEN, Inches(5.8)),
    ("CO MODULE", "DSA Operations\n& Queries", PURPLE, Inches(8.3)),
    ("SHARED\nDATA", "documents, BST,\nAVL, Graph, Costs", ORANGE, Inches(10.8)),
]

for (title, desc, color, left) in flow_steps:
    add_shape(slide, left, Inches(1.5), Inches(2.1), Inches(2.0), BG_CARD, color)
    add_text_box(slide, left + Inches(0.1), Inches(1.6), Inches(1.9), Inches(0.7), title,
                 font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.3), Inches(1.9), Inches(1.0), desc,
                 font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between boxes
for i in range(4):
    left = Inches(0.8) + Inches(2.5) * i + Inches(2.1)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, Inches(2.2), Inches(0.4), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MUTED
    arrow.line.fill.background()

# Shared data box
add_shape(slide, Inches(0.6), Inches(4.0), Inches(12.2), Inches(3.2), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.9), Inches(4.1), Inches(5), Inches(0.5), "Shared Data Architecture",
             font_size=20, color=ACCENT_LIGHT, bold=True)

shared_items = [
    "addSharedDocument(doc):",
    "  1. Check duplicate ID → add to ArrayList",
    "  2. rebuildIndexes() → reconstruct BST + AVL",
    "  3. graph.resize() → expand adjacency list",
    "  4. resizeRouteCosts() → expand cost matrix",
    "",
    "deleteSharedDocument(id):",
    "  1. Find & remove from ArrayList",
    "  2. rebuildIndexes() → reconstruct BST + AVL",
    "  3. graph.removeVertex() → remove & re-index edges",
    "  4. removeRouteCost() → shrink cost matrix"
]
add_bullet_list(slide, Inches(0.9), Inches(4.6), Inches(5.5), Inches(2.4), shared_items, font_size=12)

shared_fields = [
    "Static Fields in Main.java (single source of truth):",
    "",
    "  documents : ArrayList<Document>  → master list",
    "  bst : BST  → binary search tree index",
    "  avl : AVLTree  → self-balancing tree index",
    "  graph : Graph  → document relationship network",
    "  routeCosts : int[][]  → weighted cost matrix",
]
add_bullet_list(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.4), shared_fields, font_size=12, color=GOLD)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: Complexity Analysis
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Complexity Analysis",
             font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

complex_data = [
    ["Algorithm", "Avg Time", "Worst Time", "Space"],
    ["BST Insert/Search", "O(log n)", "O(n)", "O(n)"],
    ["AVL Insert/Search", "O(log n)", "O(log n)", "O(n)"],
    ["Fenwick Update/Query", "O(log n)", "O(log n)", "O(n)"],
    ["Segment Tree Query", "O(log n)", "O(log n)", "O(4n)"],
    ["BFS / DFS", "O(V + E)", "O(V + E)", "O(V)"],
    ["Dijkstra (matrix)", "O(V²)", "O(V²)", "O(V)"],
    ["Bellman-Ford", "O(V³)", "O(V³)", "O(V)"],
    ["Floyd-Warshall", "O(V³)", "O(V³)", "O(V²)"],
    ["Merge Sort", "O(n log n)", "O(n log n)", "O(n)"],
    ["Quick Sort", "O(n log n)", "O(n²)", "O(log n)"],
    ["Heap Sort", "O(n log n)", "O(n log n)", "O(1)"],
    ["Radix Sort", "O(d × n)", "O(d × n)", "O(n)"],
    ["0/1 Knapsack", "O(n × W)", "O(n × W)", "O(n×W)"],
    ["LIS", "O(n²)", "O(n²)", "O(n)"],
]
add_table(slide, Inches(1.5), Inches(1.3), 15, 4,
          [Inches(3), Inches(2.5), Inches(2.5), Inches(2.5)], complex_data, font_size=11)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12: How to Compile & Run
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "How to Compile & Run",
             font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# Command line
add_shape(slide, Inches(0.6), Inches(1.4), Inches(12.2), Inches(2.6), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.9), Inches(1.5), Inches(5), Inches(0.4), "Command Line",
             font_size=20, color=ACCENT_LIGHT, bold=True)
cmd_text = (
    '> cd "E:\\Apps\\Java Projects\\DocuTrack"\n'
    '> javac -d bin src/**/*.java\n'
    '> java -cp bin app.Main'
)
add_text_box(slide, Inches(0.9), Inches(2.1), Inches(11), Inches(1.5), cmd_text,
             font_size=15, color=GREEN, font_name="Consolas")

# IntelliJ
add_shape(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8), BG_CARD, PURPLE)
add_text_box(slide, Inches(0.9), Inches(4.4), Inches(5), Inches(0.4), "IntelliJ IDEA",
             font_size=20, color=PURPLE, bold=True)
ij_items = [
    '1. Open "E:\\Apps\\Java Projects" as project',
    "2. Mark DocuTrack/src as Source Root",
    "3. Right-click Main.java → Run 'Main.main()'",
    "4. Interact via the console terminal"
]
add_bullet_list(slide, Inches(0.9), Inches(4.9), Inches(5.3), Inches(2.0), ij_items, font_size=13)

# Prerequisites
add_shape(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), BG_CARD, GOLD)
add_text_box(slide, Inches(7.1), Inches(4.4), Inches(5), Inches(0.4), "Prerequisites",
             font_size=20, color=GOLD, bold=True)
prereq_items = [
    "• JDK 21 or higher",
    "• javac & java on system PATH",
    "• No external dependencies",
    "• No build tool (Maven/Gradle) needed",
    "• Pure Java SE — zero libraries"
]
add_bullet_list(slide, Inches(7.1), Inches(4.9), Inches(5.3), Inches(2.0), prereq_items, font_size=13)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13: Class Diagram
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Class Diagram",
             font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

img_path = r"E:\Apps\Java Projects\DocuTrack\class_diagram.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.3), height=Inches(6.0))
else:
    add_text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
                 "(Class diagram image not found)",
                 font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14: Thank You
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(1.5), Inches(1.8), Inches(1.8))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT
circle.line.fill.background()
add_text_box(slide, Inches(5.7), Inches(1.85), Inches(1.8), Inches(1.1), "✅", font_size=48, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(1), "Thank You!",
             font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.2), Inches(4.6), Inches(2.8), ACCENT)

add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.7),
             "DocuTrack — Where Data Structures Meet Document Management",
             font_size=20, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.1), W, Inches(0.4), BG_CARD)
add_text_box(slide, Inches(1), Inches(7.12), Inches(11), Inches(0.35),
             "6 CO Modules  •  20+ Algorithms  •  Java SE 21  •  Console Application",
             font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════
output_path = r"E:\Apps\Java Projects\DocuTrack\DocuTrack_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
