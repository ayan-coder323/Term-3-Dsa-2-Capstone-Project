from pptx import Presentation
from pptx.util import Inches

def replace_text(shape, old_text, new_text):
    if shape.has_text_frame:
        # Check if the shape's text contains the target substring
        if old_text in shape.text:
            # We want to preserve formatting, so we just replace text in paragraphs/runs
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)

def set_text(shape, target_content, new_text):
    # This just replaces the whole shape text if it matches target_content
    if shape.has_text_frame:
        if target_content in shape.text:
            shape.text = new_text

def process_presentation():
    prs = Presentation("Y25 Project Review Template - Updated.pptx")
    
    # ── SLIDE 1 ──
    slide1 = prs.slides[0]
    set_text(slide1.shapes[3], "Project Title", "DocuTrack")

    # ── SLIDE 2 ──
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            if "3-5 lines clearly describing" in shape.text:
                shape.text = "DocuTrack solves the problem of managing, indexing, and analyzing large document pools efficiently. It provides quick lookups, identifies relationships, and optimizes operations using advanced DSA, mimicking a real-world document tracking ecosystem."
            elif "Who benefits from this system?" in shape.text:
                shape.text = "Students, Administrators, and Developers learning applied DSA."
            elif "What the project aims to achieve (Goal 1)" in shape.text:
                shape.text = "Implement Core Data Structures (BST, AVL, Graph, Tries) for document management."
            elif "Key functionality or improvement (Goal 2)" in shape.text:
                shape.text = "Demonstrate Path Optimization and Storage Algorithms (Dijkstra, Knapsack, LIS)."
            elif "Final outcome or deliverable (Goal 3)" in shape.text:
                shape.text = "Deliver a robust, cohesive console application uniting all 6 CO modules."

    # ── SLIDE 3 ──
    slide3 = prs.slides[2]
    shapes_to_delete = []
    for shape in slide3.shapes:
        if shape.has_text_frame:
            replace_text(shape, "Platform: Console / Desktop", "Platform: Console Application")
            replace_text(shape, "IDE: Eclipse / IntelliJ / VS Code", "IDE: IntelliJ IDEA")
            replace_text(shape, "Classes & Objects", "Data Structures (BST, AVL, Graph)")
            replace_text(shape, "Inheritance & Polymorphism", "Graph Traversals (BFS, DFS)")
            replace_text(shape, "Interfaces", "Path Optimization (Dijkstra, Bellman-Ford)")
            replace_text(shape, "Exception Handling", "Sorting (Merge, Quick, Heap, Radix)")
            # Collection framework can remain
            replace_text(shape, "File I/O Streams", "Storage Optimization (Knapsack, LIS)")
            replace_text(shape, "Git (Version Control)", "Greedy Algorithms (Activity Selection)")
            
            if "List Major Java Concepts" in shape.text:
                shapes_to_delete.append(shape)
                
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    # ── SLIDE 4 ──
    slide4 = prs.slides[3]
    shapes_to_delete = []
    for shape in slide4.shapes:
        if shape.has_text_frame:
            replace_text(shape, "User Input", "Console Menu (Admin/User)")
            replace_text(shape, "Service Layer", "Authentication & Routing")
            replace_text(shape, "Business Logic", "DSA Operations (Algorithms)")
            replace_text(shape, "File I/O", "Central Shared Data Pool")
            replace_text(shape, "Report Generation", "Console Output Results")
            if "Draw the System Architecture" in shape.text:
                shapes_to_delete.append(shape)
                
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    # ── SLIDE 5 ──
    slide5 = prs.slides[4]
    shapes_to_delete = []
    for shape in slide5.shapes:
        if shape.has_text_frame:
            replace_text(shape, "Class 1: Primary Responsibility", "Main: Central Orchestrator & Menu System")
            replace_text(shape, "Class 2: Primary Responsibility", "Document: Core Entity with ID, Title, Author")
            replace_text(shape, "Class 3: Primary Responsibility", "Graph / BST: Shared indexing mapping")
            replace_text(shape, "Utility Class: Helper methods", "Sorting & Path Optimization classes")
            replace_text(shape, "Interface 1: Purpose & Contract", "N/A (Pure Java Implementation)")
            replace_text(shape, "Abstract Class: Shared behavior", "N/A")
            replace_text(shape, "Design Pattern: Singleton / Factory / MVC", "Pattern: Centralized Controller")
            replace_text(shape, "Usage: How it's applied in project", "Single source of truth for documents")
            if "Draw the Class Diagrams" in shape.text:
                shapes_to_delete.append(shape)
                
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    # ── SLIDE 6 ──
    slide6 = prs.slides[5]
    shapes_to_delete = []
    for shape in slide6.shapes:
        if shape.has_text_frame:
            replace_text(shape, "Start Program", "Start DocuTrack")
            replace_text(shape, "Choose Menu", "Login Admin/User")
            replace_text(shape, "Enter Input", "Select Module (1-6)")
            replace_text(shape, "Process", "Execute DSA Algorithm")
            replace_text(shape, "Output/Save", "View Computed Result")
            if "Draw the CLI /Menu Flow Diagram" in shape.text:
                shapes_to_delete.append(shape)
                
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    # ── ADD NEW SLIDE FOR CLASS DIAGRAM ──
    # Create a new blank slide (using the same layout as the current blank layout or layout 6)
    blank_slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_slide_layout)
    
    # Optional: copy background color if possible, or just insert the image
    # Add title
    txBox = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    txBox.text = "Class Diagram"
    txBox.text_frame.paragraphs[0].font.size = 360000  # Pt size equivalent approx 36
    
    img_path = r"E:\Apps\Java Projects\DocuTrack\class_diagram.png"
    import os
    if os.path.exists(img_path):
        # Insert image centrally
        new_slide.shapes.add_picture(img_path, Inches(1), Inches(1.2), width=Inches(11.3))

    # Move this new slide to position 6 (right after slide 5 which was Class Design)
    # Reordering in python-pptx requires xml manipulation
    slide_id = new_slide.slide_id
    slide_list = prs.slides._sldIdLst
    # Find the new slide element
    for el in slide_list:
        if el.id == slide_id:
            new_el = el
            break
    # Insert it after slide 5 (index 4)
    slide_list.insert(5, new_el)

    # ── SAVE ──
    prs.save("Y25 Project Review Template - Updated.pptx")
    print("Done updating template!")

process_presentation()
