from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette (Same as generate_ppt.py) ──────────────────────────
BG_DARK      = RGBColor(0x0F, 0x1B, 0x2D)   # deep navy
BG_CARD      = RGBColor(0x16, 0x2A, 0x45)   # card navy
ACCENT       = RGBColor(0x38, 0xB2, 0xAC)   # teal
ACCENT_LIGHT = RGBColor(0x81, 0xE6, 0xD9)   # light teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE0)
GOLD         = RGBColor(0xED, 0xD8, 0x8B)
GREEN        = RGBColor(0x68, 0xD3, 0x91)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def create_deep_dive_slide(prs, co_num, title_text, main_color, data_desc, logic_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # Left accent
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), main_color)
    
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.5), f"CO{co_num} Deep Dive",
                 font_size=16, color=main_color, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7), title_text,
                 font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2.5), main_color)
    
    # Data Handling Card
    add_shape(slide, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.8), BG_CARD)
    add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5), Inches(0.4), "Data Handling & Document Mapping",
                 font_size=18, color=main_color, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1), data_desc,
                 font_size=14, color=LIGHT_GRAY)
                 
    # Logic Card
    add_shape(slide, Inches(0.6), Inches(3.6), Inches(12.2), Inches(3.5), BG_CARD)
    add_text_box(slide, Inches(0.9), Inches(3.7), Inches(5), Inches(0.4), "Algorithm Logic Implementation",
                 font_size=18, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.9), Inches(4.2), Inches(11.5), Inches(2.8), logic_bullets, font_size=14)
    
    return slide

def main():
    prs = Presentation(r"E:\Apps\Java Projects\DocuTrack\DocuTrack_Presentation.pptx")
    new_slides = []
    
    # CO1
    s1 = create_deep_dive_slide(prs, 1, "Indexing Algorithm Implementation", ACCENT_LIGHT,
        "The models.Document objects are explicitly wrapped inside BSTNode and AVLNode. The algorithms route strictly on document.documentId as the primary numeric key. When addSharedDocument(doc) triggers, it cascades the insert sequentially to bst.insert(doc) and avl.insert(doc).",
        [
            "BST Logic: The insert algorithm performs an O(h) iterative walk-down, comparing document IDs, and instantiates a new BSTNode containing the exact Document reference in memory.",
            "AVL Logic: Similar to BST, but explicitly tracks subtree heights. On insertion, getBalance() checks for violations (>1 or <-1) and automatically triggers LL, RR, LR, or RL rotations to guarantee O(log n) tree shape.",
            "Retrieval: Searching traverses the numeric IDs. Upon finding a match, the algorithm returns the literal Document object reference, enabling O(1) attribute access without duplicating data."
        ])
    new_slides.append(s1)
    
    # CO2
    s2 = create_deep_dive_slide(prs, 2, "Analytics Trees Processing Logic", GOLD,
        "Unlike the persistent indexing trees, the Analytics modules (Fenwick, Segment, B-Tree) map to the central documents ArrayList dynamically. They are reconstructed on-the-fly via build() methods to guarantee data consistency.",
        [
            "Fenwick Tree (BIT): Iterates through documents using an implicit 1-indexed array. The update(i, val) function uses bitwise operations (i & -i) to add values to subsequent nodes, granting O(log n) prefix sums.",
            "Segment Tree: Allocates a 1D tree[] array of size 4N. A recursive build(node, start, end) calculates the aggregate values (e.g. sums of document IDs) for localized segments of the document pool.",
            "B+ Tree: Distributes keys across nodes with a set 'order'. The keys are explicitly documentIds. Leaf nodes retain an ArrayList<Integer> data buffer which enables direct linear rangeSearch(start, end) without recursing up the tree."
        ])
    new_slides.append(s2)

    # CO3
    s3 = create_deep_dive_slide(prs, 3, "Network Graph & Traversal Mapping", GREEN,
        "The Graph maps relationships using an adjacency list ArrayList<Integer>[] adj. Crucially, the vertices are NOT the Documents themselves, but their integer indices (0 to N-1) inside the main documents array list.",
        [
            "Dynamic Resizing: When a Document is appended globally, graph.resize(n) allocates a larger adjacency array, copying existing edges intact. Deletion triggers a complex re-indexing of all edges > deletedIndex.",
            "BFS Execution: A standard Queue pushes the starting index, exploring immediate network neighbors in O(V+E). Discovered indices are passed back to documents.get(i) to print the associated Document Title.",
            "DFS Execution: A recursive function passes a boolean[] visited array, driving deep into the network connections to map long-range document dependencies."
        ])
    new_slides.append(s3)

    # CO4
    s4 = create_deep_dive_slide(prs, 4, "Path Optimization Matrices", ACCENT,
        "The system maintains a 2D integer matrix routeCosts[u][v]. Similar to the Graph, 'u' and 'v' represent document array indices. Disconnected documents hold Dijkstra.INF (999999), and self-loops are explicitly 0.",
        [
            "Matrix Scaling: resizeRouteCosts() triggers on new document creation. It allocates a new (N+1)x(N+1) array, copies prior routing costs, and sets all new row/col bounds to INF.",
            "Dijkstra Logic: A greedy algorithm employing a boolean[] sptSet. It minimizes path costs by iterating the matrix, relaxing cost[u][v] if path(source..u) + cost(u..v) < current path(source..v).",
            "Floyd-Warshall Logic: A dynamic programming triple-nested loop (k, i, j). It evaluates if routing document 'i' through an intermediate document 'k' to reach 'j' offers a cheaper total weight."
        ])
    new_slides.append(s4)

    # CO5
    s5 = create_deep_dive_slide(prs, 5, "Sorting & Ranking Data Handling", ACCENT_LIGHT,
        "To strictly prevent mutating or corrupting the global documents list order, all sorting algorithms generate a fresh, detached deep-copy Document[] array using getDocumentArray() before executing.",
        [
            "Merge Sort: Recursively divides the array. Merges using standard numeric comparison on doc.documentId.",
            "Quick Sort: Employs Lomuto partition scheme. Pivot element logic is driven purely by doc.title.toLowerCase().compareTo() for clean alphabetical title sorting.",
            "Heap Sort: Transforms the cloned array into a Max-Heap structure using heapify(). The heap property is strictly dictated by doc.author.compareToIgnoreCase().",
            "Radix Sort: Implements counting sort recursively for every digit position based on documentId, applying (id / exp) % 10 math for robust non-comparative sorting."
        ])
    new_slides.append(s5)

    # CO6
    s6 = create_deep_dive_slide(prs, 6, "Storage DP & Greedy Processing", GOLD,
        "Storage algorithms map the generic document structures into isolated primitive integer arrays (weight[], value[], start[], finish[]) to efficiently perform intensive math without object overhead.",
        [
            "Activity Selection: Binds start/finish arrays alongside a custom Integer[] index array. Sorting indices based on finish times allows an optimal greedy extraction of non-overlapping tasks.",
            "0/1 Knapsack: Constructs a massive DP table dp[i][W]. If weight > capacity, it pulls dp[i-1][w]. Else, it calculates Math.max(val + dp[i-1][w-weight], dp[i-1][w]).",
            "Backtracking (DP Reconstruction): Instead of just returning max value, 0/1 Knapsack loops backward from dp[n][W]. When the value shifts, it identifies the exact document chosen and prints it via documents.get(i-1)."
        ])
    new_slides.append(s6)
    
    # We want to reorder these 6 slides.
    # Currently, they are at the very end of the presentation (indices -6 to -1)
    # We want to place them right before the "Execution Flow" slide.
    # Let's find the Execution Flow slide index
    exec_idx = -1
    for idx, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame and "Execution Flow" in shape.text:
                exec_idx = idx
                break
        if exec_idx != -1:
            break
            
    if exec_idx != -1:
        # Reordering in python-pptx
        slide_list = prs.slides._sldIdLst
        # The last 6 sldId elements belong to our new slides
        sIds = list(slide_list)[-6:]
        
        # Remove them from the end
        for sId in sIds:
            slide_list.remove(sId)
            
        # Insert them before exec_idx
        for sId in sIds:
            slide_list.insert(exec_idx, sId)
            exec_idx += 1

    prs.save(r"E:\Apps\Java Projects\DocuTrack\DocuTrack_Presentation.pptx")
    print("Added deep dive slides!")

if __name__ == "__main__":
    main()
