from pptx import Presentation
from pptx.util import Pt

prs = Presentation('DocuTrack_Presentation_final.pptx')

# Slide 6 (index 5) is "Class Design"
slide = prs.slides[5]

# Left column text shapes
left_text_indices = [6, 7, 8, 9]
left_icon_indices = [16, 17, 18, 19]

# Right column text shapes
right_text_indices = [12, 13, 14, 15]
right_icon_indices = [20, 21, 22, 23]

start_y_left = 2450000
start_y_right = 2470000
spacing = 750000
icon_offset = 100000

for i in range(4):
    # Left
    t_shape = slide.shapes[left_text_indices[i]]
    i_shape = slide.shapes[left_icon_indices[i]]
    
    new_y = start_y_left + i * spacing
    t_shape.top = new_y
    i_shape.top = new_y + icon_offset
    
    # Optional: adjust font size
    if t_shape.has_text_frame:
        for p in t_shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)

    # Right
    t_shape_r = slide.shapes[right_text_indices[i]]
    i_shape_r = slide.shapes[right_icon_indices[i]]
    
    new_y_r = start_y_right + i * spacing
    t_shape_r.top = new_y_r
    i_shape_r.top = new_y_r + icon_offset
    
    if t_shape_r.has_text_frame:
        for p in t_shape_r.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)

prs.save('DocuTrack_Presentation_final.pptx')
print("Overlap fixed.")
