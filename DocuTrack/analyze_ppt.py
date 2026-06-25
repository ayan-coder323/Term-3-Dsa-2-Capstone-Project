from pptx import Presentation

prs = Presentation("Y25 Project Review Template.pptx")

for i, slide in enumerate(prs.slides):
    print(f"--- SLIDE {i + 1} ---")
    print(f"Layout: {slide.slide_layout.name}")
    for j, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        text = ""
        if shape.has_text_frame:
            text = shape.text.strip().replace("\n", " | ")
            if text:
                print(f"  Shape {j} (Type: {shape_type}): {text}")
