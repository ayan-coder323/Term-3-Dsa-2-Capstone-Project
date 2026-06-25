import collections
from pptx import Presentation
import sys

def extract_text_and_images(ppt_path):
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        print(f"=== SLIDE {i+1} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"TEXT: {shape.text.replace(chr(10), ' | ')}")
            if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                print(f"IMAGE FOUND: {shape.name}")
        print("-" * 40)

if __name__ == "__main__":
    extract_text_and_images("DocuTrack_Presentation_final.pptx")
